"""Template/slide type logic: T1 pairing, T14 fill, slide deletion."""

import json
from pathlib import Path
from typing import Dict, List, Any

from pptx import Presentation
from pptx.util import Pt

from slide_generation.renderer.text import (
    _clear_text_frame,
    _fill_bullets,
    _parse_latex_format,
    _set_paragraph_formatted_text,
    _set_paragraph_no_bullet,
)


def _is_T1_textonly(s: Dict[str, Any]) -> bool:
    return (
        s.get("template_id") == "T1_TextOnly"
        and not s.get("images")
        and not s.get("tables")
        and not s.get("formulas")
    )


def pair_T1_to_T14(plan_path: str, write_back: bool = True) -> int:

    p = Path(plan_path)
    plan = json.loads(p.read_text(encoding="utf-8"))
    slides: List[Dict[str, Any]] = plan.get("slides", [])
    out: List[Dict[str, Any]] = []

    i, n, made = 0, len(slides), 0
    while i < n:
        cur = slides[i]
        if (
            i + 1 < n
            and _is_T1_textonly(cur)
            and _is_T1_textonly(slides[i + 1])
            and cur.get("section") == slides[i + 1].get("section")
        ):
            left, right = cur, slides[i + 1]
            out.append({
                "section": cur.get("section"),
                "template_id": "T14_2Text",
                "columns": [
                    {
                        "subsection": left.get("subsection", "") or "",
                        "bullets": left.get("bullets", []) or [],
                        "paragraph": left.get("paragraph") or ""
                    },
                    {
                        "subsection": right.get("subsection", "") or "",
                        "bullets": right.get("bullets", []) or [],
                        "paragraph": right.get("paragraph") or ""
                    }
                ]
            })
            made += 1
            i += 2
            continue

        out.append(cur)
        i += 1

    plan["slides"] = out
    if write_back:
        p.write_text(json.dumps(plan, ensure_ascii=False, indent=4), encoding="utf-8")
    return made


def validate_no_consecutive_T1(plan_path: str) -> List[int]:
    p = Path(plan_path)
    plan = json.loads(p.read_text(encoding="utf-8"))
    slides: List[Dict[str, Any]] = plan.get("slides", [])
    bad_idxs: List[int] = []
    for i in range(len(slides) - 1):
        a, b = slides[i], slides[i + 1]
        if _is_T1_textonly(a) and _is_T1_textonly(b) and a.get("section") == b.get("section"):
            bad_idxs.append(i)
    return bad_idxs


def _get_placeholder(slide, name):
    for shp in slide.shapes:
        if getattr(shp, "name", "").strip().lower() == name.strip().lower():
            return shp
    return None


def _ph_text_n(slide, n: int):
    targets = {f"text placeholder {n}", f"\u6587\u672c\u5360\u4f4d\u7b26 {n}"}
    for shp in slide.shapes:
        nm = getattr(shp, "name", "").strip().lower()
        if nm in targets:
            return shp
        if ("placeholder" in nm or "\u5360\u4f4d\u7b26" in nm) and nm.endswith(f" {n}"):
            return shp
    return None


def _ph_by_idx(slide, idx: int):
    for ph in getattr(slide, "placeholders", []):
        if ph.placeholder_format.idx == idx:
            return ph
    return None


def fill_T14_2Text(slide, slide_info, section_no_text, use_paragraph: bool = True, bullet_font_size: int = 20, title_font_size: int = 28):
    print("slide_info:")
    print(slide_info)
    part_ph = (
        _get_placeholder(slide, "Part")
        or _get_placeholder(slide, "Text Placeholder 2")
        or _ph_text_n(slide, 2)
    )

    title_bar = _get_placeholder(slide, "Text Placeholder 1") or _ph_text_n(slide, 1)
    title_bar = _ph_by_idx(slide, 1)  or _ph_text_n(slide, 1)


    section_title = slide_info.get("section")

    title_bar.text = section_title
    tf = title_bar.text_frame
    tf.clear()
    tf.paragraphs[0].text = section_title
    print("[AFTER] title_bar text =", repr(tf.text))

    lt = (
        _get_placeholder(slide, "Left Title")
        or _get_placeholder(slide, "Text Placeholder 3")
        or _ph_text_n(slide, 3)
    )
    lb = (
        _get_placeholder(slide, "Left Body")
        or _get_placeholder(slide, "Text Placeholder 4")
        or _ph_text_n(slide, 4)
    )
    rt = (
        _get_placeholder(slide, "Right Title")
        or _get_placeholder(slide, "Text Placeholder 5")
        or _ph_text_n(slide, 5)
    )
    rb = (
        _get_placeholder(slide, "Right Body")
        or _get_placeholder(slide, "Text Placeholder 6")
        or _ph_text_n(slide, 6)
    )


    cols = slide_info.get("columns") or []
    left  = cols[0] if len(cols) > 0 else {}
    right = cols[1] if len(cols) > 1 else {}

    if part_ph is not None and getattr(part_ph, "has_text_frame", False):
        part_ph.text_frame.text = f"{section_no_text}"

    # Ensure T14_2Text titles use sizes consistent with other slides
    if title_bar is not None and getattr(title_bar, "has_text_frame", False):
        title_txt = slide_info.get("section", "") or slide_info.get("title", "")
        ttf = title_bar.text_frame
        ttf.text = title_txt
        for p in ttf.paragraphs:
            p.font.size = Pt(title_font_size)

    if lt is not None and getattr(lt, "has_text_frame", False):
        ltf = lt.text_frame
        ltf.text = left.get("subsection", "") or left.get("title", "") or ""
        for p in ltf.paragraphs:
            p.font.size = Pt(bullet_font_size)

    if rt is not None and getattr(rt, "has_text_frame", False):
        rtf = rt.text_frame
        rtf.text = right.get("subsection", "") or right.get("title", "") or ""
        for p in rtf.paragraphs:
            p.font.size = Pt(bullet_font_size)

    # Use same bullet/paragraph logic as other layouts: prefer bullets when present; use paragraph when bullets empty
    def _fill_body_tf(tf, col: dict):
        bullets = col.get("bullets") or []
        paragraph_text = (col.get("paragraph") or "").strip()
        if use_paragraph and paragraph_text and not bullets:
            _clear_text_frame(tf)
            if tf.paragraphs:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.level = 0
            _set_paragraph_no_bullet(p)
            segs = _parse_latex_format(paragraph_text)
            if segs:
                _set_paragraph_formatted_text(p, segs, font_size=bullet_font_size)
            else:
                p.text = paragraph_text
                p.font.size = Pt(bullet_font_size)
        else:
            _fill_bullets(tf, bullets, lvl0_size=bullet_font_size, lvl1_size=bullet_font_size)

    if lb is not None and getattr(lb, "has_text_frame", False):
        _fill_body_tf(lb.text_frame, left)

    if rb is not None and getattr(rb, "has_text_frame", False):
        _fill_body_tf(rb.text_frame, right)


    missing = []
    for name, ph in [
        ("Part(2)", part_ph), ("Left Title(3)", lt), ("Left Body(4)", lb),
        ("Right Title(5)", rt), ("Right Body(6)", rb),
    ]:
        if ph is None:
            missing.append(name)
    if missing:
        print(f"[WARN] T14_2Text \u6a21\u677f\u7f3a\u5c11\u5360\u4f4d\uff1a{', '.join(missing)}")


def delete_slide(prs: Presentation, slide_index: int) -> None:
    """
    Delete slide by index in python-pptx (works by removing sldId and dropping rel).
    """
    sldIdLst = prs.slides._sldIdLst  # pylint: disable=protected-access
    sldId_elems = list(sldIdLst)

    if slide_index < 0 or slide_index >= len(sldId_elems):
        raise IndexError(f"slide_index out of range: {slide_index}, total={len(sldId_elems)}")

    sldId = sldId_elems[slide_index]
    rId = sldId.rId
    sldIdLst.remove(sldId)
    prs.part.drop_rel(rId)
