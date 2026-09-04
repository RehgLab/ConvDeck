"""slide_generation.renderer -- slide rendering package.

Re-exports the public API so that ``from slide_generation.renderer import *``
(or explicit imports) work for the renderer's public API.
"""

from slide_generation.renderer.visuals import (  # noqa: F401
    _insert_picture_keep_ratio,
    insert_image_below_content,
    insert_visuals_auto,
    resolve_visual_paths,
    resolve_formula_mode1_path,
    verify_and_fix_plan_visuals,
    get_image_reasons,
    get_table_reasons,
    get_formula_reasons,
    _best_match,
    _collect_reasons_for_kind,
    _extract_idx,
    _nums_from_files,
)

from slide_generation.renderer.text import (  # noqa: F401
    LATEX_COLOR_MAP,
    _bullet_word_count,
    _bullet_line_count,
    _find_matching_brace,
    _parse_latex_format,
    _color_to_rgb_tuple,
    _set_run_color,
    _set_paragraph_formatted_text,
    _set_paragraph_no_bullet,
    _clear_text_frame,
    _set_text_frame_columns,
    _fill_bullets,
    set_font_color,
)

from slide_generation.renderer.slide_types import (  # noqa: F401
    _is_T1_textonly,
    pair_T1_to_T14,
    validate_no_consecutive_T1,
    fill_T14_2Text,
    delete_slide,
    _get_placeholder,
    _ph_text_n,
    _ph_by_idx,
)

from slide_generation.renderer.placeholders import (  # noqa: F401
    TEXT_TYPES,
    find_text_placeholders,
    get_content,
    _placeholder_by_name,
    scan_layout_placeholders,
    extract_theme_color_from_title,
)

# Shared module-level color constants
from pptx.dml.color import RGBColor as _RGBColor  # noqa: F401

COLOR_WHITE = _RGBColor(0, 0, 0)
THEME_COLOR = _RGBColor(185, 210, 153)
