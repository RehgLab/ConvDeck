"""Placeholder discovery functions."""

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.shapes import PP_PLACEHOLDER_TYPE as PH_TYPE

from slide_generation.renderer.visuals import _best_match


TEXT_TYPES = {
    PH_TYPE.TITLE,
    PH_TYPE.CENTER_TITLE,
    PH_TYPE.SUBTITLE,
    PH_TYPE.BODY,
}


def find_text_placeholders(slide):
    """Return (part_num_ph, subsection_ph, body_ph) by position, or None if insufficient placeholders."""
    txt_ph = [
        s for s in slide.shapes
        if (
            s.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER and
            s.has_text_frame and
            s.placeholder_format.type in TEXT_TYPES
        )
    ]

    if len(txt_ph) < 3:
        print(f"Warning: Slide has only {len(txt_ph)} text placeholders, skipping text content...")
        return None


    # Robust selection:
    # - Identify BODY by placeholder type (avoid relying on the "3rd by top" assumption).
    body_candidates = [s for s in txt_ph if s.placeholder_format.type == PH_TYPE.BODY]
    if not body_candidates:
        print("Warning: No BODY placeholder found on slide, skipping text content...")
        return None

    # Body is typically the lowest text region on the slide.
    body_ph = max(body_candidates, key=lambda s: s.top)

    upper_candidates = [s for s in txt_ph if s is not body_ph]
    if len(upper_candidates) < 2:
        print("Warning: Insufficient non-BODY placeholders on slide, skipping text content...")
        return None

    upper_candidates_sorted = sorted(upper_candidates, key=lambda s: (s.top, s.left))
    top_two = sorted(upper_candidates_sorted[:2], key=lambda s: s.left)
    part_ph, title_ph = top_two

    return part_ph, title_ph, body_ph


def get_content(sec_title, sub_title, outline):
    _, sub = _best_match(outline, sec_title, sub_title)
    return sub.get("content", "") if sub else ""


def _placeholder_by_name(slide, name: str):
    """Return placeholder shape whose .name == name."""

    for shape in slide.shapes:
        if shape.name == name:
            return shape
    raise KeyError(f'Placeholder "{name}" not found on slide master.')


def scan_layout_placeholders(template_path: str):
    prs = Presentation(template_path)
    layout_map = {}

    for layout in prs.slide_layouts:
        layout_name = layout.name
        placeholder_names = []

        for shape in layout.shapes:
            shape_type = shape.shape_type
            shape_type_name = str(shape_type)
            try:
                shape_type_name = MSO_SHAPE_TYPE(shape_type).name
            except ValueError:
                pass

            placeholder_names.append({
                "name": shape.name,
                "type": shape_type_name
            })

        layout_map[layout_name] = placeholder_names

    return layout_map


def extract_theme_color_from_title(prs, layout_index=2):
    slide_layout = prs.slide_layouts[layout_index]
    temp_slide = prs.slides.add_slide(slide_layout)

    for shape in temp_slide.shapes:
        if (
            shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER and
            shape.placeholder_format.type == 1  # 1 == TITLE
        ):
            para = shape.text_frame.paragraphs[0]
            font_color = para.font.color
            if font_color.type == 2:  # 2 == THEME
                theme_color = font_color.theme_color
                xml_slides = prs.slides._sldIdLst
                xml_slides.remove(xml_slides[-1])  # remove last slide
                return theme_color
    return None
