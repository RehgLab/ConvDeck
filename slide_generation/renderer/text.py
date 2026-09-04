"""Text/bullet formatting functions and color-parsing helpers."""

from typing import List

from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.oxml.xmlchemy import OxmlElement
from pptx.oxml.ns import qn


# LaTeX-like formatting: \textbf{...} and \textcolor{<color>}{...}
LATEX_COLOR_MAP = {
    "red": (255, 0, 0),
    "blue": (0, 0, 255),
    "green": (0, 128, 0),
    "black": (0, 0, 0),
    "white": (255, 255, 255),
    "orange": (255, 165, 0),
    "yellow": (255, 255, 0),
    "purple": (128, 0, 128),
    "brown": (165, 42, 42),
    "gray": (128, 128, 128),
    "grey": (128, 128, 128),
}


def _bullet_word_count(bullets):
    """Return total word count across all bullet text and sub-bullets."""
    if not bullets:
        return 0
    n = 0
    for b in bullets:
        n += len((b.get("text") or "").split())
        for s in b.get("sub") or []:
            n += len(str(s).split())
    return n


def _bullet_line_count(bullets):
    """Return total number of lines (each main bullet + each sub-bullet counts as one line)."""
    if not bullets:
        return 0
    return sum(1 + len(b.get("sub") or []) for b in bullets)


def _find_matching_brace(s: str, start: int) -> int:
    """Return index of closing } for the { at start. start must point to {."""
    depth = 1
    i = start + 1
    while i < len(s) and depth > 0:
        if s[i] == "{":
            depth += 1
        elif s[i] == "}":
            depth -= 1
        i += 1
    return i - 1 if depth == 0 else -1


def _parse_latex_format(s: str) -> List[tuple]:
    """
    Parse string for \\textbf{...} and \\textcolor{<color>}{...}.
    Also matches textbf{/extbf{ and textcolor{/extcolor{ in case backslash was stripped or \\t became tab.
    Returns list of (text, bold, color) where color is RGBColor or None.
    """
    if not s:
        return []
    # Normalize: \\t in JSON/strings can become a tab, so \\textbf{ -> tab + "extbf{"; accept both
    PATTERNS_BF = ("\\textbf{", "textbf{", "extbf{")
    PATTERNS_TC = ("\\textcolor{", "textcolor{", "extcolor{")
    out = []
    i = 0
    while i < len(s):
        # Find next match and which pattern matched (for correct brace offset)
        idx_bf, len_bf = len(s), 0
        for pat in PATTERNS_BF:
            pos = s.find(pat, i)
            if pos != -1 and pos < idx_bf:
                idx_bf, len_bf = pos, len(pat)
        idx_tc, len_tc = len(s), 0
        for pat in PATTERNS_TC:
            pos = s.find(pat, i)
            if pos != -1 and pos < idx_tc:
                idx_tc, len_tc = pos, len(pat)
        next_idx = min(idx_bf, idx_tc)
        if next_idx > i:
            plain = s[i:next_idx]
            # Remove trailing tabs (and skip segment if only whitespace) to avoid extra tab/space before LaTeX-formatted words
            plain = plain.rstrip("\t")
            if plain.strip():
                out.append((plain, False, None))
            i = next_idx
            continue
        if next_idx == len(s):
            break
        if idx_bf <= idx_tc:
            # \textbf{ or textbf{ or extbf{ -> open_br is the { (at len_bf - 1 from start)
            open_br = next_idx + len_bf - 1
            close_br = _find_matching_brace(s, open_br)
            if close_br == -1:
                i = next_idx + 1
                continue
            content = s[open_br + 1 : close_br]
            for txt, bold, color in _parse_latex_format(content):
                out.append((txt, True, color))
            i = close_br + 1
        else:
            # \textcolor{ or textcolor{ or extcolor{
            open_br = next_idx + len_tc - 1  # index of first {
            close_br = _find_matching_brace(s, open_br)
            if close_br == -1:
                i = next_idx + 1
                continue
            color_name = s[open_br + 1 : close_br].strip().lower()
            next_open = close_br + 1
            if next_open < len(s) and s[next_open] == "{":
                content_close = _find_matching_brace(s, next_open)
                if content_close != -1:
                    content = s[next_open + 1 : content_close]
                    rgb = LATEX_COLOR_MAP.get(color_name)
                    use_color = RGBColor(*rgb) if rgb else None
                    for txt, bold, color in _parse_latex_format(content):
                        out.append((txt, bold, use_color or color))
                    i = content_close + 1
                else:
                    out.append((s[next_idx:next_open + 1], False, None))
                    i = next_open + 1
            else:
                out.append((s[next_idx:close_br + 1], False, None))
                i = close_br + 1
    return out


def _color_to_rgb_tuple(color):
    """Return (r, g, b) from RGBColor or (r,g,b) tuple."""
    if color is None:
        return None
    if isinstance(color, (list, tuple)) and len(color) == 3:
        return (int(color[0]), int(color[1]), int(color[2]))
    rgb = getattr(color, "_rgb", None)
    if rgb is not None:
        return ((rgb >> 16) & 0xFF, (rgb >> 8) & 0xFF, rgb & 0xFF)
    return (getattr(color, "r", 0), getattr(color, "g", 0), getattr(color, "b", 0))


def _set_run_color(run, color) -> None:
    """Set run font color; when run.font.color is None, set via XML."""
    rgb_tuple = _color_to_rgb_tuple(color)
    if rgb_tuple is None:
        return
    try:
        if getattr(run.font, "color", None) is not None:
            run.font.color.rgb = RGBColor(*rgb_tuple)
            return
    except (AttributeError, TypeError):
        pass
    # Fallback: set color via run's XML (run.font.color is None in some layouts)
    try:
        rPr = run.font._element
        for child in list(rPr.iterchildren()):
            if child.tag == qn("a:solidFill"):
                rPr.remove(child)
        solid_fill = OxmlElement("a:solidFill")
        srgb = OxmlElement("a:srgbClr")
        r, g, b = rgb_tuple
        srgb.set("val", f"{r:02X}{g:02X}{b:02X}")
        solid_fill.append(srgb)
        rPr.append(solid_fill)
    except Exception:
        pass


def _set_paragraph_formatted_text(paragraph, segments: List[tuple], font_size: int = 20):
    """Set paragraph content from (text, bold, color) segments using runs."""
    segments = [(t, b, c) for t, b, c in segments if t]
    if not segments:
        paragraph.text = ""
        return
    # Clear existing runs (paragraph.text = "" may leave zero runs in some pptx layouts)
    for run in list(paragraph.runs):
        run._element.getparent().remove(run._element)
    # Add one run per segment
    for text, bold, color in segments:
        run = paragraph.add_run()
        run.text = text
        run.font.size = Pt(font_size)
        run.font.bold = bold
        _set_run_color(run, color)


def _set_paragraph_no_bullet(paragraph):
    """Turn off bullet and remove indent so paragraph displays as plain left-aligned text."""
    p_el = getattr(paragraph, "_p", getattr(paragraph, "_element", None))
    if p_el is None:
        return
    pPr = p_el.find(qn("a:pPr"))
    if pPr is None:
        pPr = OxmlElement("a:pPr")
        p_el.insert(0, pPr)
    for tag in ("a:buChar", "a:buBlip", "a:buAutoNum", "a:buFont", "a:buSzPct"):
        existing = pPr.find(qn(tag))
        if existing is not None:
            pPr.remove(existing)
    bu_none = OxmlElement("a:buNone")
    pPr.insert(0, bu_none)
    # Remove left margin and indent so all lines align left (no tab/indent on wrap)
    pPr.set("marL", "0")
    pPr.set("indent", "0")
    # Remove level-based indent if present (lvl)
    if pPr.get("lvl") is not None:
        pPr.attrib.pop("lvl", None)
    # Also clear indent via API so paragraph_format is consistent
    try:
        pf = paragraph.paragraph_format
        pf.left_indent = Pt(0)
        pf.first_line_indent = Pt(0)
    except Exception:
        pass


def _clear_text_frame(tf):
    if tf.paragraphs:
        tf.paragraphs[0].text = ""

        while len(tf.paragraphs) > 1:
            p = tf._element.p_lst[-1]
            p.getparent().remove(p)
    else:
        tf.clear()


def _set_text_frame_columns(tf, num_cols):
    """Set number of columns for a text frame via underlying XML (e.g. 2 for double column)."""
    for c in tf._element.getchildren():
        if c.tag.split("}")[-1] == "bodyPr":
            c.set("numCol", str(num_cols))
            return
    # Fallback: try setting on first bodyPr found by local name
    try:
        from pptx.oxml.ns import qn as _qn
        body_pr = tf._element.find(_qn("a:bodyPr"))
        if body_pr is not None:
            body_pr.set("numCol", str(num_cols))
    except Exception:
        pass


def _fill_bullets(tf, bullets, lvl0_size=20, lvl1_size=20):
    _clear_text_frame(tf)
    for b in (bullets or []):
        p = tf.add_paragraph()
        p.level = 0
        raw_text = (b.get("text") or "").strip()
        segments = _parse_latex_format(raw_text)
        if segments:
            _set_paragraph_formatted_text(p, segments, font_size=lvl0_size)
        else:
            p.text = raw_text
            p.font.size = Pt(lvl0_size)

        for s in (b.get("sub") or []):
            sp = tf.add_paragraph()
            sp.level = 1
            raw_sub = str(s).strip()
            segs_sub = _parse_latex_format(raw_sub)
            if segs_sub:
                _set_paragraph_formatted_text(sp, segs_sub, font_size=lvl1_size)
            else:
                sp.text = raw_sub
                sp.font.size = Pt(lvl1_size)


def set_font_color(paragraph, theme_color):
    paragraph.font.fill.solid()
    if theme_color is not None:
        paragraph.font.fill.fore_color.theme_color = theme_color
    else:
        paragraph.font.fill.fore_color.rgb = RGBColor(255, 105, 180)  # pink fallback
