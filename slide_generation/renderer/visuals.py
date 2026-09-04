"""Image/visual insertion functions and helpers."""

import json
import re
from pathlib import Path

import difflib
from PIL import Image
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Pt


def _insert_picture_keep_ratio(ph, img_path: Path):
    slide = ph.part.slide
    ph_left, ph_top = ph.left, ph.top
    ph_width, ph_height = ph.width, ph.height

    with Image.open(img_path) as img:
        iw, ih = img.size
    aspect = iw / ih

    if ph_width / ph_height > aspect:
        new_h = ph_height
        new_w = int(new_h * aspect)
    else:
        new_w = ph_width
        new_h = int(new_w / aspect)

    new_left = ph_left + int((ph_width - new_w) / 2)
    new_top = ph_top + int((ph_height - new_h) / 2)

    pic = slide.shapes.add_picture(str(img_path), new_left, new_top, width=new_w, height=new_h)

    # --- move picture *behind* all text placeholders ---
    spTree = slide.shapes._spTree
    spTree.remove(pic.element)         # temporarily take it out
    spTree.insert(2, pic.element)      # index 0=background,1=layout

    # --- finally, remove the now-unused picture placeholder itself ---
    ph.element.getparent().remove(ph.element)


def insert_image_below_content(slide, img_path: Path):
    """
    Insert an image below the lowest existing shape (text or image) on the slide.
    Centered horizontally. Resizes if not enough space.
    """

    with Image.open(img_path) as img:
        width_px, height_px = img.size
    aspect_ratio = width_px / height_px

    # slide size (EMUs)
    slide_width = slide.part.slide_layout.part.package.presentation_part.slide_width
    slide_height = slide.part.slide_layout.part.package.presentation_part.slide_height


    target_width = slide_width * 0.6
    target_height = target_width / aspect_ratio


    lowest_bottom = 0
    for shape in slide.shapes:
        bottom = shape.top + shape.height
        if bottom > lowest_bottom:
            lowest_bottom = bottom

    margin = Pt(20)
    available_space = slide_height - lowest_bottom - margin

    if available_space < target_height:
        target_height = available_space
        target_width = target_height * aspect_ratio
        if target_height <= 0:
            print("Not enough space to insert image:", img_path)
            return


    left = (slide_width - target_width) // 2
    top = lowest_bottom + margin

    slide.shapes.add_picture(str(img_path), left, top, width=int(target_width), height=int(target_height))


def insert_visuals_auto(slide, visuals: list[Path]):
    """
    Automatically insert visuals (images, tables, formulas) into all available
    picture placeholders on the given slide.
    """
    # Find all picture placeholder shapes
    picture_placeholders = [
        shape for shape in slide.shapes
        if shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER and
           "Picture" in shape.name
    ]

    if len(visuals) > len(picture_placeholders):
        print(f"Warning: not enough picture placeholders on slide (needed {len(visuals)}, found {len(picture_placeholders)})")

    #  Insert images one by one
    for img_path, ph in zip(visuals, picture_placeholders):
        _insert_picture_keep_ratio(ph, Path(img_path))


def _extract_idx(val):
    if isinstance(val, int):
        return val
    m = re.findall(r"\d+", str(val))
    return int(m[-1]) if m else None


def _nums_from_files(files):
    """['image_7.png', 'table_1.png'] -> {7, 1}"""
    out = set()
    for f in files:
        m = re.findall(r"\d+", str(f))
        if m:
            out.add(int(m[-1]))
    return out


def _best_match(data, sec_title: str, sub_title: str, min_ratio: float = 0.55):
    best_sec, best_sec_score = None, 0.0
    for sec in data.get("sections", []):
        s = difflib.SequenceMatcher(None, sec.get("title","").lower(), (sec_title or "").lower()).ratio()
        if s > best_sec_score:
            best_sec, best_sec_score = sec, s
    if not best_sec or best_sec_score < min_ratio:
        return None, None

    best_sub, best_sub_score = None, 0.0
    for sub in best_sec.get("subsections", []):
        s = difflib.SequenceMatcher(None, sub.get("title","").lower(), (sub_title or "").lower()).ratio()
        if s > best_sub_score:
            best_sub, best_sub_score = sub, s
    if not best_sub or best_sub_score < min_ratio:
        return None, None
    return best_sec, best_sub


def _collect_reasons_for_kind(sec, sub, files, figs_data, kind: str):

    _, sub_d = _best_match(figs_data, sec, sub)
    if not sub_d:
        return ""

    want = _nums_from_files(files)
    if not want:
        return ""

    #   N -> value_idx
    pairs = []  # [(N, value_idx)]
    for k, v in sub_d.items():
        k_low = str(k).lower()
        if not k_low.startswith(kind):  # 'image' or 'table'
            continue
        N = _extract_idx(k)
        v_idx = _extract_idx(v)
        if N is None or v_idx is None:
            continue
        if v_idx in want:
            pairs.append((N, v_idx))

    reasons = []

    only_one_asset = (len(pairs) == 1)

    for N, _v_idx in pairs:

        candidate_keys = [
            f"reason{N}",
        ]
        if kind == "table":
            candidate_keys.append(f"reasonT{N}")
        else:
            candidate_keys.extend([f"reasonI{N}", f"reasonImg{N}"])

        if only_one_asset:
            candidate_keys.append("reason")


        found = None
        for rk in candidate_keys:
            if rk in sub_d and isinstance(sub_d[rk], str) and sub_d[rk].strip():
                found = sub_d[rk].strip()
                break


        if not found and "reason" in sub_d and isinstance(sub_d["reason"], str) and sub_d["reason"].strip():
            found = sub_d["reason"].strip()

        if found:
            reasons.append(found)

    return "\n".join(reasons)


def get_image_reasons(sec, sub, image_files, figs_data):
    return _collect_reasons_for_kind(sec, sub, image_files, figs_data, kind="image")


def get_table_reasons(sec, sub, table_files, figs_data):
    return _collect_reasons_for_kind(sec, sub, table_files, figs_data, kind="table")


def get_formula_reasons(sec, sub, formula_files, formula_data) -> str:

    _, sub_d = _best_match(formula_data, sec, sub)
    if not sub_d:
        return ""

    pairs = []
    for k, v in sub_d.items():
        k_low = str(k).lower()
        if not k_low.startswith("formula"):
            continue
        N = _extract_idx(k)
        if N is None:
            continue
        rkey = f"reason{N}"
        rtxt = sub_d.get(rkey)
        if isinstance(rtxt, str) and rtxt.strip():
            pairs.append((N, rtxt.strip()))


    try:
        pairs.sort(key=lambda x: int(x[0]))
    except Exception:
        pass


    if not pairs and isinstance(sub_d.get("reason"), str) and sub_d["reason"].strip():
        if formula_files and len(formula_files) == 1:
            return sub_d["reason"].strip()

    if formula_files and len(pairs) > len(formula_files):
        pairs = pairs[:len(formula_files)]

    out, seen = [], set()
    for _, r in pairs:
        if r not in seen:
            out.append(r)
            seen.add(r)

    return "\n".join(out)


def resolve_formula_mode1_path(fname: str, args) -> Path:
    """
    Extract formula index i from fname (like "formula_4.png"),
    and generate the path:
    <{args.model_name_t}_{args.model_name_v}>_images_and_tables/{args.paper_name}/{args.paper_name}-formula-i.png
    """

    stem = Path(fname).stem
    match = re.search(r'(\d+)(?!.*\d)', stem)
    if not match:
        raise ValueError(f"Cannot extract index from formula filename: {fname}")

    i = match.group(1)
    path_str = f"<{args.model_name_t}_{args.model_name_v}>_images_and_tables/{args.paper_name}/{args.paper_name}-formula-{i}.png"
    return Path(path_str)


def verify_and_fix_plan_visuals(args, plan_path):
    """
    One-pass sanity check over the slide plan: for every entry in each slide's
    ``images`` / ``tables`` list, look up the filename in the paper's
    ``*_images.json`` and ``*_tables.json``. If an entry sits in the wrong
    list (e.g. a table file in ``images``), swap it. Writes the corrected
    plan back to disk.
    """

    paper = args.paper_name
    prefix = f"<{args.model_name_t}_{args.model_name_v}>_images_and_tables"

    images_json_path = Path(prefix) / f"{paper}_images.json"
    tables_json_path = Path(prefix) / f"{paper}_tables.json"
    images_json = json.load(open(images_json_path, "r", encoding="utf-8")) if images_json_path.exists() else {}
    tables_json = json.load(open(tables_json_path, "r", encoding="utf-8")) if tables_json_path.exists() else {}

    image_files = {Path(v.get("image_path", "")).name for v in images_json.values()}
    table_files = {Path(v.get("table_path", "")).name for v in tables_json.values()}

    plan_path = Path(plan_path)
    with open(plan_path, "r", encoding="utf-8") as f:
        plan = json.load(f)

    swaps = 0
    for slide in plan.get("slides", []):
        new_images, new_tables = [], list(slide.get("tables", []) or [])
        for entry in slide.get("images", []) or []:
            name = Path(str(entry)).name
            if name in table_files and name not in image_files:
                new_tables.append(entry)
                swaps += 1
                print(f"[visuals-check] moved {name!r} images -> tables")
            else:
                new_images.append(entry)

        fixed_tables = []
        for entry in new_tables:
            name = Path(str(entry)).name
            if name in image_files and name not in table_files:
                new_images.append(entry)
                swaps += 1
                print(f"[visuals-check] moved {name!r} tables -> images")
            else:
                fixed_tables.append(entry)

        slide["images"] = new_images
        slide["tables"] = fixed_tables

    if swaps:
        with open(plan_path, "w", encoding="utf-8") as f:
            json.dump(plan, f, indent=4, ensure_ascii=False)
        print(f"[visuals-check] Fixed {swaps} misplaced visuals in {plan_path}")


def resolve_visual_paths(slide_info, args):


    paper = args.paper_name
    prefix = f"<{args.model_name_t}_{args.model_name_v}>_images_and_tables"
    base_dir = Path(prefix) / paper

    images_json_path = Path(prefix) / f"{paper}_images.json"
    tables_json_path = Path(prefix) / f"{paper}_tables.json"
    images_json = json.load(open(images_json_path, "r", encoding="utf-8")) if images_json_path.exists() else {}
    tables_json = json.load(open(tables_json_path, "r", encoding="utf-8")) if tables_json_path.exists() else {}

    if args.formula_mode == 3:
        formulas_dir = Path("contents") / paper / "formula_images"
    else:
        formulas_dir = base_dir

    def _norm_digits(s: str) -> str:

        m = re.search(r'(\d+)(?!.*\d)', str(s))
        if not m:
            return str(s)
        v = m.group(1).lstrip("0")
        return v if v != "" else "0"

    def _stem_name(p: str) -> str:
        return Path(str(p)).name

    def _extract_figure_numbers(caption: str):

        out = set()
        if not caption:
            return out
        caps = str(caption)

        #  Figure 6-7 / Fig. 10-12
        for m in re.finditer(r'(?:(?:fig(?:ure)?|图)\.?\s*)(\d+)\s*[–-]\s*(\d+)', caps, flags=re.I):
            a, b = int(m.group(1)), int(m.group(2))
            if a <= b:
                out.update(range(a, b + 1))
            else:
                out.update(range(b, a + 1))

        #  Figure 6 / Fig. 6
        for m in re.finditer(r'(?:(?:fig(?:ure)?|图)\.?\s*)(\d+)', caps, flags=re.I):
            out.add(int(m.group(1)))

        return out

    def _extract_table_numbers(caption: str):

        out = set()
        if not caption:
            return out
        caps = str(caption)

        for m in re.finditer(r'(?:(?:tab(?:le)?|表)\.?\s*)(\d+)\s*[–-]\s*(\d+)', caps, flags=re.I):
            a, b = int(m.group(1)), int(m.group(2))
            if a <= b:
                out.update(range(a, b + 1))
            else:
                out.update(range(b, a + 1))

        for m in re.finditer(r'(?:(?:tab(?:le)?|表)\.?\s*)(\d+)', caps, flags=re.I):
            out.add(int(m.group(1)))

        return out

    def _build_caption_index_images(images_dict: dict):

        idx = {}
        for k, meta in (images_dict or {}).items():
            caps = (meta or {}).get("caption", "")
            for f in _extract_figure_numbers(caps):
                idx.setdefault(f, set()).add(str(k))
        return idx

    def _build_caption_index_tables(tables_dict: dict):

        idx = {}
        for k, meta in (tables_dict or {}).items():
            caps = (meta or {}).get("caption", "")
            for t in _extract_table_numbers(caps):
                idx.setdefault(t, set()).add(str(k))
        return idx

    fignum_to_imgkeys = _build_caption_index_images(images_json)
    tbnum_to_tbkeys  = _build_caption_index_tables(tables_json)


    images_json = {str(k): v for k, v in images_json.items()}
    tables_json = {str(k): v for k, v in tables_json.items()}

    def _resolve_and_check(path1: Path) -> Path:

        if path1.exists():
            return path1

        p2 = base_dir / path1.name
        if p2.exists():
            return p2

        candidates = [p for p in base_dir.glob("*") if p.name.lower() == path1.name.lower()]
        if candidates:
            return candidates[0]

        existing = "\n".join(sorted(p.name for p in base_dir.glob("*.*"))[:80])
        raise FileNotFoundError(
            f"Missing visual file: {path1}\n"
            f"Tried: {path1}, {p2}. Base dir: {base_dir}\n"
            f"Existing (first 80):\n{existing}"
        )

    def _match_by_id_or_caption(img_id: str, name: str, mapping: dict, cap_index: dict, kind: str):

        img_id_norm = _norm_digits(img_id)
        name_norm   = _norm_digits(_stem_name(name))

        if img_id in mapping:
            return mapping[img_id], img_id
        if img_id_norm in mapping:
            return mapping[img_id_norm], img_id_norm
        if name_norm in mapping:
            return mapping[name_norm], name_norm

        if img_id_norm.isdigit():
            num = int(img_id_norm)
            cand = sorted(cap_index.get(num, []), key=lambda x: int(_norm_digits(x)) if _norm_digits(x).isdigit() else 10**9)
            if cand:
                k = cand[0]
                return mapping[k], k

        if name_norm.isdigit():
            num2 = int(name_norm)
            cand2 = sorted(cap_index.get(num2, []), key=lambda x: int(_norm_digits(x)) if _norm_digits(x).isdigit() else 10**9)
            if cand2:
                k2 = cand2[0]
                return mapping[k2], k2

        avail = sorted(mapping.keys(), key=lambda x: int(_norm_digits(x)) if _norm_digits(x).isdigit() else 10**9)
        print(
            f"[resolve_visual_paths] WARN: {kind} id not found, dropping. "
            f"given_id='{img_id}', file='{name}', "
            f"norm(id)={img_id_norm}, norm(file)={name_norm}. "
            f"Available keys (first 40): {avail[:40]}"
        )
        return None, None

    # ----------------- images -----------------
    image_paths = []
    for img_str in slide_info.get("images", []):
        name = _stem_name(img_str)

        #   'picture|image|fig|table|formula-<num>.ext'
        m = re.search(r'(?:picture|image|fig|table|formula)[-_](\d+)(?=\.[A-Za-z0-9]+$)', name, re.I)
        img_id = m.group(1) if m else _norm_digits(name)


        rec, used_key = _match_by_id_or_caption(img_id, name, images_json, fignum_to_imgkeys, kind="image")
        if rec is None:
            continue

        target = rec.get("image_path") or rec.get("path") or rec.get("file")
        if not target:
            raise KeyError(f"Image record has no path field. key={used_key}, record={rec}")

        img_path = _resolve_and_check(Path(target))
        image_paths.append(img_path)

    # ----------------- tables -----------------
    table_paths = []
    for tb_str in slide_info.get("tables", []):
        name = _stem_name(tb_str)
        m = re.search(r'(?:table|tab|picture|image|fig|formula)[-_](\d+)(?=\.[A-Za-z0-9]+$)', name, re.I)
        tb_id = m.group(1) if m else _norm_digits(name)

        rec, used_key = _match_by_id_or_caption(tb_id, name, tables_json, tbnum_to_tbkeys, kind="table")
        if rec is None:
            continue

        target = rec.get("table_path") or rec.get("image_path") or rec.get("path") or rec.get("file")
        if not target:
            raise KeyError(f"Table record has no path field. key={used_key}, record={rec}")

        tb_path = _resolve_and_check(Path(target))
        table_paths.append(tb_path)

    # ----------------- formulas -----------------
    formula_paths = []
    for fname in slide_info.get("formulas", []):
        if args.formula_mode == 3:
            final_path = formulas_dir / fname
        else:
            final_path = Path(resolve_formula_mode1_path(fname, args))
        formula_paths.append(_resolve_and_check(Path(final_path)))
    print("formula_paths",formula_paths)
    return image_paths + table_paths + formula_paths
