"""Visual helpers: rendering, borders, labels, and image utilities."""

import asyncio
from pathlib import Path
from urllib.parse import quote_from_bytes, quote

from PIL import Image
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from playwright.sync_api import sync_playwright
from playwright.async_api import async_playwright


def add_border_to_all_elements(prs, border_color=RGBColor(255, 0, 0), border_width=Pt(2)):
    """
    Iterates over all slides and shapes in the Presentation object 'prs'
    and applies a red border with the specified width to each shape.

    Args:
        prs: The Presentation object to modify.
        border_color: An instance of RGBColor for the border color (default is red).
        border_width: The width of the border as a Pt value (default is 2 points).
    """
    for slide in prs.slides:
        for shape in slide.shapes:
            # Some shapes (like charts or group shapes) might not support border styling
            try:
                # Set the line fill to be solid and assign the desired color and width.
                shape.line.fill.solid()
                shape.line.fill.fore_color.rgb = border_color
                shape.line.width = border_width
            except Exception as e:
                # If a shape doesn't support setting a border, print a message and continue.
                print(f"Could not add border to shape {shape.shape_type}: {e}")


# 1 point = 12700 EMUs (helper function)
def pt_to_emu(points: float) -> int:
    return int(points * 12700)


def add_border_and_labels(
    prs,
    border_color=RGBColor(255, 0, 0),   # Red border for shapes
    border_width=Pt(2),                # 2-point border width
    label_outline_color=RGBColor(0, 0, 255),  # Blue outline for label circle
    label_text_color=RGBColor(0, 0, 255),     # Blue text color
    label_diameter_pt=40                       # Diameter of the label circle in points
):
    """
    Iterates over all slides and shapes in the Presentation 'prs', applies a
    red border to each shape, and places a transparent (no fill), blue-outlined
    circular label with a blue number in the center of each shape. Labels start
    from 0 and increment for every shape that gets a border.

    Args:
        prs: The Presentation object to modify.
        border_color: RGBColor for the shape border color (default: red).
        border_width: The width of the shape border (Pt).
        label_outline_color: The outline color for the label circle (default: blue).
        label_text_color: The color of the label text (default: blue).
        label_diameter_pt: The diameter of the label circle, in points (default: 40).
    """
    label_diameter_emu = pt_to_emu(label_diameter_pt)  # convert diameter (points) to EMUs
    label_counter = 0  # Start labeling at 0
    labeled_elements = {}

    for slide in prs.slides:
        for shape in slide.shapes:
            # Skip shapes that are labels themselves
            if shape.name.startswith("Label_"):
                continue

            try:
                # --- 1) Add red border to the shape (if supported) ---
                shape.line.fill.solid()
                shape.line.fill.fore_color.rgb = border_color
                shape.line.width = border_width

                # --- 2) Calculate center for the label circle ---
                label_left = shape.left + (shape.width // 2) - (label_diameter_emu // 2)
                label_top  = shape.top  + (shape.height // 2) - (label_diameter_emu // 2)

                # --- 3) Create label circle (an OVAL) in the center of the shape ---
                label_shape = slide.shapes.add_shape(
                    MSO_AUTO_SHAPE_TYPE.OVAL,
                    label_left,
                    label_top,
                    label_diameter_emu,
                    label_diameter_emu
                )
                label_shape.name = f"Label_{label_counter}"  # so we can skip it later

                # **Make the circle completely transparent** (no fill at all)
                label_shape.fill.background()

                # **Give it a blue outline**
                label_shape.line.fill.solid()
                label_shape.line.fill.fore_color.rgb = label_outline_color
                label_shape.line.width = Pt(3)

                # --- 4) Add the label number (centered, blue text) ---
                tf = label_shape.text_frame
                tf.text = str(label_counter)
                paragraph = tf.paragraphs[0]
                paragraph.alignment = PP_ALIGN.CENTER

                run = paragraph.runs[0]
                font = run.font
                font.size = Pt(40)      # Larger font
                font.bold = True
                font.name = "Arial"
                font._element.get_or_change_to_solidFill()
                font.fill.fore_color.rgb = label_text_color
                # Record properties from the original shape and label text.
                labeled_elements[label_counter] = {
                    'left': f'{shape.left} EMU',
                    'top': f'{shape.top} EMU',
                    'width': f'{shape.width} EMU',
                    'height': f'{shape.height} EMU',
                    'font_size': f'{shape.text_frame.font.size} PT' if hasattr(shape, 'text_frame') else None,
                }

                # --- 5) Increment label counter (so every shape has a unique label) ---
                label_counter += 1

            except Exception as e:
                # If the shape doesn't support borders or text, skip gracefully
                print(f"Could not add border/label to shape (type={shape.shape_type}): {e}")

    return labeled_elements


async def rendered_dims(html: Path) -> tuple[int, int]:
    async with async_playwright() as p:
        browser = await p.chromium.launch()
        page    = await browser.new_page()        # no fixed viewport yet
        resolved = html.resolve()
        # quote_from_bytes expects bytes, so we encode the path as UTF-8:
        url = "file://" + quote_from_bytes(str(resolved).encode("utf-8"), safe="/:")
        await page.goto(url, wait_until="networkidle")

        # 1) bounding-box of <body>
        body_box = await page.eval_on_selector(
            "body",
            "el => el.getBoundingClientRect()")
        w = int(body_box["width"])
        h = int(body_box["height"])

        await browser.close()
        return w, h


def html_to_png(html_abs_path, slides_width_default, slides_height_default, output_path):
    html_file = html_abs_path

    try:
        w, h = asyncio.run(rendered_dims(html_file))
    except:
        w = slides_width_default
        h = slides_height_default

    with sync_playwright() as p:
        path_posix = Path(html_file).resolve().as_posix()

        file_url = "file://" + quote(path_posix, safe="/:")
        browser = p.chromium.launch()
        page    = browser.new_page(viewport={"width": w, "height": h})
        page.goto(file_url, wait_until='networkidle')
        page.screenshot(path=output_path, full_page=True)
        browser.close()


def get_img_ratio(img_path):
    img = Image.open(img_path)
    return {
        'width': img.width,
        'height': img.height
    }


def get_img_ratio_in_section(content_json):
    res = {}
    if 'path' in content_json:
        res[content_json['path']] = get_img_ratio(content_json['path'])

    if 'subsections' in content_json:
        for subsection_name, val in content_json['subsections'].items():
            if 'path' in val:
                res[val['path']] = get_img_ratio(val['path'])

    return res


def remove_hierarchy_and_id(data):
    """
    Recursively remove the 'hierarchy' and 'id' fields from a nested
    dictionary representing sections and subsections.
    """
    if isinstance(data, dict):
        # Create a new dict to store filtered data
        new_data = {}
        for key, value in data.items():
            # Skip the keys "hierarchy" and "id"
            if key in ("hierarchy", "id", 'location'):
                continue
            # Recursively process the value
            new_data[key] = remove_hierarchy_and_id(value)
        return new_data
    elif isinstance(data, list):
        # If it's a list, process each item recursively
        return [remove_hierarchy_and_id(item) for item in data]
    else:
        # Base case: if it's neither dict nor list, just return the value as is
        return data
